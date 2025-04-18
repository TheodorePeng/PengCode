import os
import glob
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import tempfile

def merge_ppt_files(folder_paths):
    current_time = datetime.now().strftime('%Y%m%d_%H%M%S')
    output_file = os.path.join(folder_paths[0], f'{current_time}.pptx')
    
    merged_presentation = Presentation()
    
    for folder_path in folder_paths:
        print(f'Processing folder: {folder_path}')
        ppt_files = sorted(glob.glob(os.path.join(folder_path, '*.pptx')))
        
        for ppt_file in ppt_files:
            presentation = Presentation(ppt_file)
            
            for slide in presentation.slides:
                slide_layout = merged_presentation.slide_layouts[0]
                new_slide = merged_presentation.slides.add_slide(slide_layout)
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_box = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                        text_frame = text_box.text_frame
                        text_frame.text = shape.text_frame.text
                    elif shape.shape_type == 13:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_img_file:
                            tmp_img_file.write(shape.image.blob)
                            tmp_img_file.flush()
                            new_slide.shapes.add_picture(tmp_img_file.name, shape.left, shape.top, shape.width, shape.height)
                        os.remove(tmp_img_file.name)

    merged_presentation.save(output_file)
    print(f'Merged presentation saved as: {output_file}')

if __name__ == "__main__":
    # Initial run
    folder_paths = ['/Users/peng/Downloads/11']
    merge_ppt_files(folder_paths)
    
    # Manually change folder_paths and run again
    folder_paths = ['/Users/peng/Downloads/12']
    merge_ppt_files(folder_paths)