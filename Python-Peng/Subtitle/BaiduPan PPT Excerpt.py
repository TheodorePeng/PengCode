from pptx import Presentation
from PIL import Image
import os

# 定义PPT文件路径
ppt_file = '/Users/peng/Downloads/11/20240624_100654.pptx'  # 自定义设置PPT文件路径

# 创建一个Presentation对象
prs = Presentation(ppt_file)

# 获取当前用户的桌面路径，并创建一个新文件夹来保存提取的图片
desktop_path = os.path.join(os.path.expanduser('~'), 'Desktop')
image_folder = os.path.join(desktop_path, 'PPT images')
if not os.path.exists(image_folder):
    os.makedirs(image_folder)

# 遍历每一张幻灯片并提取图片
for slide_num, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # 13 corresponds to Picture
            image = shape.image
            image_bytes = image.blob
            image_filename = os.path.join(image_folder, f'image_{slide_num + 1}.png')
            with open(image_filename, 'wb') as img_file:
                img_file.write(image_bytes)
            
            # 确保图片保存为PNG格式
            img = Image.open(image_filename)
            img = img.convert('RGB')  # 确保图片是RGB格式
            img.save(image_filename, 'PNG')

print(f"图片已保存至：{image_folder}")

# 需要自定义设置的参数:
# ppt_file = '/path/to/your/ppt_file.pptx'  # 设置PPT文件路径

# Desc. 百度网盘 提取视频截图 PPT
    # PPT特点：这个PPT中每张slide中都只放置了一张图片
    # 你只需要将PPT文件路径改成你自己的路径，然后运行脚本即可。
    # 脚本会自动遍历每一张幻灯片，并提取图片，保存到桌面。图片会自动保存为PNG格式。图片命名规则为：image_1.png, image_2.png, image_3.png...